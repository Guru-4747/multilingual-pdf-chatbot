import os
import io
import re
import csv
import base64
import hashlib

import faiss
import numpy as np
import streamlit as st

from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer
from google import genai
from google.genai import types

# Multi-format / multimodal extraction
import fitz  # PyMuPDF - handles PDF text, page rasterization (for scans), and embedded images

try:
    import docx  # python-docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation  # python-pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# =========================================================
# CONFIGURATION
# =========================================================

st.set_page_config(
    page_title="Multilingual AI Document Chatbot",
    page_icon="📚",
    layout="wide"
)

GENERATION_MODEL = "gemini-2.5-flash"

SUPPORTED_EXTENSIONS = [
    "pdf", "docx", "pptx",
    "txt", "csv",
    "png", "jpg", "jpeg", "webp"
]

IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "webp"}
IMAGE_MIME = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "webp": "image/webp",
}


# =========================================================
# CREATIVE / COLOURFUL UI
# =========================================================

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');
html, body, [class*="css"] { font-family: "Inter", sans-serif; }
.stApp { background: radial-gradient(circle at 8% 8%, rgba(124,58,237,.30), transparent 28%), radial-gradient(circle at 92% 10%, rgba(14,165,233,.22), transparent 28%), radial-gradient(circle at 50% 95%, rgba(236,72,153,.18), transparent 32%), #070b18; }
.block-container { max-width: 1250px; padding-top: 2rem; padding-bottom: 4rem; }
.hero { padding: 34px 36px; border-radius: 28px; background: linear-gradient(135deg, rgba(124,58,237,.34), rgba(37,99,235,.22), rgba(236,72,153,.20)); border: 1px solid rgba(255,255,255,.12); box-shadow: 0 24px 70px rgba(0,0,0,.35); margin-bottom: 24px; }
.hero-title { font-size: clamp(32px, 5vw, 50px); font-weight: 800; color: white; line-height: 1.05; margin: 12px 0 10px; }
.hero-subtitle { color: #dbeafe; font-size: 18px; line-height: 1.6; }
.online { display: inline-block; padding: 7px 13px; border-radius: 999px; background: rgba(34,197,94,.13); color: #86efac; border: 1px solid rgba(34,197,94,.30); font-size: 12px; font-weight: 700; }
.card { background: rgba(15,23,42,.72); backdrop-filter: blur(16px); border: 1px solid rgba(255,255,255,.09); border-radius: 22px; padding: 24px; margin: 16px 0; box-shadow: 0 14px 40px rgba(0,0,0,.22); }
.card-title { color: white; font-size: 20px; font-weight: 750; margin-bottom: 8px; }
.card-text { color: #94a3b8; line-height: 1.6; }
.upload-intro { text-align: center; padding: 28px 20px 12px; }
.upload-icon { font-size: 54px; }
.upload-title { color: white; font-size: 27px; font-weight: 800; margin-top: 8px; }
.upload-text { color: #94a3b8; margin-top: 8px; }
[data-testid="stFileUploaderDropzone"] { background: rgba(15,23,42,.60); border: 2px dashed rgba(139,92,246,.60); border-radius: 20px; }
.stat-card { min-height: 128px; padding: 20px 14px; border-radius: 20px; text-align: center; background: linear-gradient(145deg, rgba(30,41,59,.86), rgba(15,23,42,.72)); border: 1px solid rgba(255,255,255,.08); box-shadow: 0 10px 30px rgba(0,0,0,.20); }
.stat-icon { font-size: 28px; }
.stat-value { color: white; font-size: 25px; font-weight: 800; margin-top: 5px; }
.stat-label { color: #94a3b8; font-size: 12px; margin-top: 3px; }
.lang-card { padding: 15px 8px; text-align: center; border-radius: 16px; background: linear-gradient(135deg, rgba(124,58,237,.17), rgba(59,130,246,.10)); border: 1px solid rgba(255,255,255,.07); color: #f8fafc; font-weight: 650; margin-bottom: 10px; }
.section-heading { color: white; font-size: 23px; font-weight: 800; margin: 26px 0 10px; }
.feature-card { padding: 18px; min-height: 125px; border-radius: 18px; background: rgba(15,23,42,.66); border: 1px solid rgba(255,255,255,.07); }
.feature-title { color: white; font-weight: 750; font-size: 16px; }
.feature-text { color: #94a3b8; font-size: 13px; line-height: 1.5; margin-top: 7px; }
.source-card { padding: 13px 16px; border-radius: 14px; background: rgba(124,58,237,.10); border: 1px solid rgba(124,58,237,.20); color: #c4b5fd; margin-top: 8px; }
.type-badge { display:inline-block; font-size: 11px; font-weight: 700; padding: 3px 9px; border-radius: 999px; margin-right: 6px; }
.type-text { background: rgba(59,130,246,.18); color: #93c5fd; border: 1px solid rgba(59,130,246,.30); }
.type-image { background: rgba(236,72,153,.18); color: #f9a8d4; border: 1px solid rgba(236,72,153,.30); }
.type-ocr { background: rgba(234,179,8,.18); color: #fde68a; border: 1px solid rgba(234,179,8,.30); }
.stButton > button { border-radius: 14px; font-weight: 700; border: 1px solid rgba(255,255,255,.08); transition: all .2s ease; }
.stButton > button:hover { transform: translateY(-2px); box-shadow: 0 10px 25px rgba(124,58,237,.25); }
section[data-testid="stSidebar"] { background: linear-gradient(180deg, #0b1020, #111827); border-right: 1px solid rgba(255,255,255,.08); }
.footer { text-align: center; color: #64748b; padding: 34px 10px 10px; font-size: 12px; line-height: 1.8; }
</style>
""", unsafe_allow_html=True)

load_dotenv()


# =========================================================
# API KEY
# =========================================================

api_key = os.getenv("GEMINI_API_KEY")

if not api_key:
    st.error("Gemini API key not found.")
    st.info("Create a .env file and add GEMINI_API_KEY=YOUR_KEY")
    st.stop()

client = genai.Client(api_key=api_key)


# =========================================================
# MULTILINGUAL EMBEDDING MODEL
# =========================================================

@st.cache_resource
def load_embedding_model():
    return SentenceTransformer(
        "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
    )


embedding_model = load_embedding_model()


# =========================================================
# LANGUAGES
# =========================================================

LANGUAGES = {
    "English": "English",
    "தமிழ் (Tamil)": "Tamil",
    "ಕನ್ನಡ (Kannada)": "Kannada",
    "हिन्दी (Hindi)": "Hindi",
    "తెలుగు (Telugu)": "Telugu",
    "മലയാളం (Malayalam)": "Malayalam"
}


# =========================================================
# SESSION STATE
# =========================================================

for key, default in [
    ("processed", False),
    ("chunks", []),
    ("index", None),
    ("file_hash", None),
    ("messages", []),
]:
    if key not in st.session_state:
        st.session_state[key] = default


# =========================================================
# TITLE
# =========================================================

st.markdown("""
<div class="hero">
    <div class="online">🟢 AI SYSTEM ONLINE</div>
    <div class="hero-title">📚 LinguaDoc AI</div>
    <div class="hero-subtitle">Your intelligent multilingual, multimodal document assistant</div>
    <div style="color:#cbd5e1; margin-top:12px; line-height:1.6;">
        Upload a PDF, Word doc, PowerPoint, text file, or even a photo. LinguaDoc AI reads
        the text <b>and</b> the images inside your document (charts, diagrams, scanned pages,
        photos) using Gemini's vision, then answers your questions, summarizes, and
        generates MCQs in English, Tamil, Kannada, Hindi, Telugu or Malayalam.
    </div>
</div>
""", unsafe_allow_html=True)


# =========================================================
# SIDEBAR
# =========================================================

with st.sidebar:

    st.header("⚙️ Settings")

    selected_language = st.selectbox(
        "Answer Language",
        list(LANGUAGES.keys())
    )

    top_k = st.slider(
        "Relevant chunks to retrieve",
        min_value=2,
        max_value=12,
        value=6
    )

    describe_images = st.checkbox(
        "Analyze images inside the document (Gemini Vision)",
        value=True,
        help="Extracts embedded images/diagrams/charts and scanned pages, and asks "
             "Gemini to describe/OCR them so they become searchable too. Slower but "
             "much more accurate for image-heavy documents."
    )

    strict_grounding = st.checkbox(
        "Strict grounding (only answer from document)",
        value=True,
        help="If off, Gemini may lightly use general knowledge to clarify terms, "
             "while still prioritizing the document."
    )

    st.divider()

    st.write("### Supported File Types")
    st.write("📄 PDF · 📝 DOCX · 📊 PPTX")
    st.write("📃 TXT / CSV · 🖼️ PNG / JPG / WEBP")

    st.divider()

    st.write("### Supported Languages")
    st.write("🇬🇧 English")
    st.write("🇮🇳 தமிழ்")
    st.write("🇮🇳 ಕನ್ನಡ")
    st.write("🇮🇳 हिन्दी")
    st.write("🇮🇳 తెలుగు")
    st.write("🇮🇳 മലയാളം")


# =========================================================
# UPLOADER
# =========================================================

st.markdown("""
<div class="card">
    <div class="upload-intro">
        <div class="upload-icon">📄</div>
        <div class="upload-title">Upload Your Document</div>
        <div class="upload-text">PDF, Word, PowerPoint, text, or an image. Scanned pages and
        embedded pictures are read too, not just selectable text.</div>
    </div>
</div>
""", unsafe_allow_html=True)

uploaded_file = st.file_uploader(
    "Choose a file",
    type=SUPPORTED_EXTENSIONS,
    label_visibility="collapsed",
)


# =========================================================
# HELPERS
# =========================================================

def get_file_hash(uploaded_file):
    file_bytes = uploaded_file.getvalue()
    return hashlib.md5(file_bytes).hexdigest()


def get_extension(filename):
    return filename.rsplit(".", 1)[-1].lower()


def clean_text(text):
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def create_chunks(text, page_number, source_type="text", chunk_size=700, overlap=100):
    """Split text into overlapping word-based chunks, tagged with page and source type."""
    words = text.split()
    chunks = []
    start = 0

    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk_text = " ".join(words[start:end])

        if chunk_text.strip():
            chunks.append({
                "text": chunk_text,
                "page": page_number,
                "type": source_type,
            })

        if end == len(words):
            break
        start = end - overlap

    return chunks


# =========================================================
# GEMINI VISION HELPERS (image description + OCR)
# =========================================================

def describe_image_with_gemini(image_bytes, mime_type="image/png", context_hint=""):
    """Ask Gemini to describe an image AND transcribe any visible text (OCR)."""
    try:
        prompt = (
            "Carefully analyze this image, which was extracted from a document"
            f"{(' ' + context_hint) if context_hint else ''}. "
            "1) Transcribe any visible text exactly (OCR). "
            "2) Describe charts, diagrams, tables, photos, or figures in detail, "
            "including axis labels, numbers, and key relationships if present. "
            "Be factual and specific — do not guess at content you cannot see clearly. "
            "Respond in plain text only."
        )

        import os
import streamlit as st
from google import genai
from google.genai import types

# ==============================
# GEMINI API KEY
# ==============================

try:
    api_key = st.secrets["GEMINI_API_KEY"]
except Exception:
    api_key = os.getenv("GEMINI_API_KEY")

if not api_key:
    st.error("Gemini API key not found.")
    st.info(
        "Please add GEMINI_API_KEY in "
        "Streamlit Cloud → Manage app → Settings → Secrets."
    )
    st.stop()

client = genai.Client(api_key=api_key)

GENERATION_MODEL = "gemini-2.5-flash"

def ocr_scanned_page_with_gemini(image_bytes, page_number):
    """Full-page OCR for scanned / image-only PDF pages via Gemini vision."""
    try:
        prompt = (
            f"This is page {page_number} of a scanned document with no selectable text. "
            "Transcribe ALL visible text exactly as written, preserving reading order, "
            "including headers, tables, and captions. If there is no readable text, say "
            "'NO_TEXT_FOUND'. Respond with the transcription only, no commentary."
        )

        response = client.models.generate_content(
            model=GENERATION_MODEL,
            contents=[
                types.Part.from_bytes(data=image_bytes, mime_type="image/png"),
                prompt,
            ],
        )

        text = clean_text(response.text or "")
        if "NO_TEXT_FOUND" in text:
            return ""
        return text

    except Exception:
        return ""


# =========================================================
# PDF EXTRACTION (text + scanned pages + embedded images)
# =========================================================

def extract_pdf_chunks(uploaded_file, analyze_images):
    file_bytes = uploaded_file.getvalue()
    pdf = fitz.open(stream=file_bytes, filetype="pdf")

    total_pages = pdf.page_count
    chunks = []

    progress = st.progress(0)
    status = st.empty()

    for page_index in range(total_pages):
        page_number = page_index + 1
        page = pdf[page_index]

        status.write(f"Reading page {page_number} of {total_pages}...")

        # --- 1. Extract selectable text ---
        try:
            raw_text = page.get_text("text") or ""
        except Exception:
            raw_text = ""

        text = clean_text(raw_text)

        if text:
            chunks.extend(create_chunks(text, page_number, source_type="text"))
        elif analyze_images:
            # --- 2. No selectable text -> likely a scanned page. Rasterize + OCR. ---
            status.write(f"Page {page_number} looks scanned — running OCR via Gemini...")
            pix = page.get_pixmap(dpi=200)
            png_bytes = pix.tobytes("png")
            ocr_text = ocr_scanned_page_with_gemini(png_bytes, page_number)
            if ocr_text:
                chunks.extend(create_chunks(ocr_text, page_number, source_type="ocr"))

        # --- 3. Extract embedded images (diagrams, charts, photos) on this page ---
        if analyze_images:
            try:
                image_list = page.get_images(full=True)
            except Exception:
                image_list = []

            for img_index, img in enumerate(image_list, start=1):
                try:
                    xref = img[0]
                    base_image = pdf.extract_image(xref)
                    img_bytes = base_image["image"]
                    img_ext = base_image.get("ext", "png")

                    # Skip tiny decorative images (icons, bullets, lines)
                    if len(img_bytes) < 3000:
                        continue

                    mime_type = f"image/{img_ext}" if img_ext != "jpg" else "image/jpeg"
                    status.write(
                        f"Analyzing image {img_index} on page {page_number} with Gemini Vision..."
                    )
                    description = describe_image_with_gemini(
                        img_bytes, mime_type=mime_type,
                        context_hint=f"on page {page_number}"
                    )
                    if description:
                        chunks.extend(
                            create_chunks(
                                f"Image on page {page_number}: {description}",
                                page_number,
                                source_type="image",
                            )
                        )
                except Exception:
                    continue

        progress.progress(page_number / total_pages)

    status.success(f"Finished reading {total_pages} pages.")
    progress.empty()
    pdf.close()

    return chunks, total_pages


# =========================================================
# DOCX EXTRACTION (paragraphs, tables, embedded images)
# =========================================================

def extract_docx_chunks(uploaded_file, analyze_images):
    if not DOCX_AVAILABLE:
        st.error("python-docx is not installed. Run: pip install python-docx")
        st.stop()

    document = docx.Document(io.BytesIO(uploaded_file.getvalue()))

    status = st.empty()
    status.write("Reading Word document...")

    full_text = []
    for para in document.paragraphs:
        if para.text.strip():
            full_text.append(para.text)

    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                full_text.append(row_text)

    text = clean_text("\n".join(full_text))
    chunks = create_chunks(text, page_number="Document", source_type="text") if text else []

    if analyze_images:
        try:
            image_parts = [
                rel.target_part for rel in document.part.rels.values()
                if "image" in rel.reltype
            ]
        except Exception:
            image_parts = []

        for idx, part in enumerate(image_parts, start=1):
            try:
                img_bytes = part.blob
                if len(img_bytes) < 3000:
                    continue
                mime_type = part.content_type or "image/png"
                status.write(f"Analyzing image {idx} in document with Gemini Vision...")
                description = describe_image_with_gemini(
                    img_bytes, mime_type=mime_type, context_hint="in a Word document"
                )
                if description:
                    chunks.extend(
                        create_chunks(
                            f"Image {idx} in document: {description}",
                            page_number="Document",
                            source_type="image",
                        )
                    )
            except Exception:
                continue

    status.success("Finished reading the Word document.")
    return chunks, 1


# =========================================================
# PPTX EXTRACTION (slide text + embedded images)
# =========================================================

def extract_pptx_chunks(uploaded_file, analyze_images):
    if not PPTX_AVAILABLE:
        st.error("python-pptx is not installed. Run: pip install python-pptx")
        st.stop()

    presentation = Presentation(io.BytesIO(uploaded_file.getvalue()))

    total_slides = len(presentation.slides)
    chunks = []

    progress = st.progress(0)
    status = st.empty()

    for slide_index, slide in enumerate(presentation.slides, start=1):
        status.write(f"Reading slide {slide_index} of {total_slides}...")

        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        slide_text.append(line)

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_text.append(row_text)

        text = clean_text("\n".join(slide_text))
        if text:
            chunks.extend(create_chunks(text, page_number=f"Slide {slide_index}", source_type="text"))

        if analyze_images:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    try:
                        img_bytes = shape.image.blob
                        if len(img_bytes) < 3000:
                            continue
                        mime_type = shape.image.content_type or "image/png"
                        status.write(
                            f"Analyzing image on slide {slide_index} with Gemini Vision..."
                        )
                        description = describe_image_with_gemini(
                            img_bytes, mime_type=mime_type,
                            context_hint=f"on slide {slide_index}"
                        )
                        if description:
                            chunks.extend(
                                create_chunks(
                                    f"Image on slide {slide_index}: {description}",
                                    page_number=f"Slide {slide_index}",
                                    source_type="image",
                                )
                            )
                    except Exception:
                        continue

        progress.progress(slide_index / total_slides)

    status.success(f"Finished reading {total_slides} slides.")
    progress.empty()

    return chunks, total_slides


# =========================================================
# TXT / CSV EXTRACTION
# =========================================================

def extract_text_file_chunks(uploaded_file):
    raw_bytes = uploaded_file.getvalue()

    try:
        raw_text = raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        raw_text = raw_bytes.decode("latin-1", errors="ignore")

    ext = get_extension(uploaded_file.name)

    if ext == "csv":
        rows = list(csv.reader(io.StringIO(raw_text)))
        lines = [" | ".join(row) for row in rows]
        raw_text = "\n".join(lines)

    text = clean_text(raw_text)
    chunks = create_chunks(text, page_number="Document", source_type="text") if text else []

    return chunks, 1


# =========================================================
# STANDALONE IMAGE EXTRACTION
# =========================================================

def extract_image_file_chunks(uploaded_file):
    ext = get_extension(uploaded_file.name)
    mime_type = IMAGE_MIME.get(ext, "image/png")
    img_bytes = uploaded_file.getvalue()

    status = st.empty()
    status.write("Analyzing image with Gemini Vision...")

    description = describe_image_with_gemini(
        img_bytes, mime_type=mime_type, context_hint="uploaded directly by the user"
    )

    status.success("Finished analyzing the image.")

    chunks = create_chunks(description, page_number="Image", source_type="image") if description else []
    return chunks, 1


# =========================================================
# UNIFIED DISPATCHER
# =========================================================

def extract_document_chunks(uploaded_file, analyze_images):
    ext = get_extension(uploaded_file.name)

    if ext == "pdf":
        return extract_pdf_chunks(uploaded_file, analyze_images)
    elif ext == "docx":
        return extract_docx_chunks(uploaded_file, analyze_images)
    elif ext == "pptx":
        return extract_pptx_chunks(uploaded_file, analyze_images)
    elif ext in ("txt", "csv"):
        return extract_text_file_chunks(uploaded_file)
    elif ext in IMAGE_EXTENSIONS:
        return extract_image_file_chunks(uploaded_file)
    else:
        st.error(f"Unsupported file type: .{ext}")
        st.stop()


# =========================================================
# EMBEDDINGS + FAISS INDEX
# =========================================================

def create_embeddings(chunks):
    texts = [item["text"] for item in chunks]
    embeddings = []
    batch_size = 32

    progress = st.progress(0)
    total = len(texts)

    for start in range(0, total, batch_size):
        batch = texts[start:start + batch_size]
        batch_embeddings = embedding_model.encode(
            batch, normalize_embeddings=True, show_progress_bar=False
        )
        embeddings.extend(batch_embeddings)
        progress.progress(min((start + batch_size) / total, 1.0))

    progress.empty()
    return np.asarray(embeddings, dtype="float32")


def create_faiss_index(embeddings):
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatIP(dimension)
    index.add(embeddings)
    return index


# =========================================================
# PROCESS UPLOADED FILE
# =========================================================

if uploaded_file:

    current_hash = get_file_hash(uploaded_file)
    file_changed = current_hash != st.session_state.file_hash

    if file_changed:
        st.session_state.processed = False
        st.session_state.chunks = []
        st.session_state.index = None
        st.session_state.messages = []
        st.session_state.file_hash = current_hash

    ext = get_extension(uploaded_file.name)
    icon = "🖼️" if ext in IMAGE_EXTENSIONS else "📎"

    st.markdown(f"""
<div class="card">
    <div class="card-title">{icon} Selected Document</div>
    <div class="card-text">{uploaded_file.name}</div>
</div>
""", unsafe_allow_html=True)

    if st.button("✨ Analyze My Document", type="primary", use_container_width=True):

        with st.spinner("Processing document. Image-heavy or scanned files may take longer..."):

            chunks, unit_count = extract_document_chunks(uploaded_file, describe_images)

            if not chunks:
                st.error("No readable content was found.")
                st.warning(
                    "Try enabling 'Analyze images inside the document' in the sidebar if this "
                    "is a scanned PDF or an image-only file."
                )
                st.stop()

            unit_label = "slides" if ext == "pptx" else "pages" if ext == "pdf" else "sections"
            st.success(
                f"📚 Processed {unit_count} {unit_label} and created "
                f"{len(chunks):,} searchable chunks (text + image-derived)."
            )

            embeddings = create_embeddings(chunks)
            index = create_faiss_index(embeddings)

            st.session_state.chunks = chunks
            st.session_state.index = index
            st.session_state.processed = True

            st.success("🎉 Your document is ready for questions!")

            n_image_chunks = sum(1 for c in chunks if c["type"] in ("image", "ocr"))

            c1, c2, c3, c4 = st.columns(4)
            with c1:
                st.markdown(f'<div class="stat-card"><div class="stat-icon">📄</div><div class="stat-value">{unit_count}</div><div class="stat-label">{unit_label.title()}</div></div>', unsafe_allow_html=True)
            with c2:
                st.markdown(f'<div class="stat-card"><div class="stat-icon">🧩</div><div class="stat-value">{len(chunks):,}</div><div class="stat-label">Text Chunks</div></div>', unsafe_allow_html=True)
            with c3:
                st.markdown(f'<div class="stat-card"><div class="stat-icon">🖼️</div><div class="stat-value">{n_image_chunks}</div><div class="stat-label">Image/OCR Chunks</div></div>', unsafe_allow_html=True)
            with c4:
                st.markdown('<div class="stat-card"><div class="stat-icon">🌐</div><div class="stat-value">6</div><div class="stat-label">Languages</div></div>', unsafe_allow_html=True)


# =========================================================
# LANGUAGE SUPPORT DISPLAY
# =========================================================

st.markdown('<div class="section-heading">🌐 Language Support</div>', unsafe_allow_html=True)

lang_cols = st.columns(6)
for col, language in zip(lang_cols, ["🇬🇧 English", "🇮🇳 தமிழ்", "🇮🇳 ಕನ್ನಡ", "🇮🇳 हिन्दी", "🇮🇳 తెలుగు", "🇮🇳 മലയാళം"]):
    with col:
        st.markdown(f'<div class="lang-card">{language}</div>', unsafe_allow_html=True)


# =========================================================
# SEARCH
# =========================================================

def search_document(question, chunks, index, top_k):
    question_embedding = embedding_model.encode([question], normalize_embeddings=True)
    question_embedding = np.asarray(question_embedding, dtype="float32")

    scores, positions = index.search(question_embedding, min(top_k, len(chunks)))

    results = []
    for score, position in zip(scores[0], positions[0]):
        if position == -1:
            continue
        results.append({
            "text": chunks[position]["text"],
            "page": chunks[position]["page"],
            "type": chunks[position]["type"],
            "score": float(score),
        })

    return results


# =========================================================
# GEMINI ANSWER GENERATION
# =========================================================

def generate_answer(question, results, language, strict=True):
    context_parts = []

    for result in results:
        label = {
            "text": "Text",
            "image": "Image description",
            "ocr": "Scanned page OCR",
        }.get(result["type"], "Content")

        context_parts.append(
            f"[{label} — Location: {result['page']}]\n{result['text']}"
        )

    context = "\n\n".join(context_parts)

    grounding_rule = (
        "Answer the user's question using ONLY the information provided in the "
        "DOCUMENT CONTEXT below. Do not use outside knowledge."
        if strict else
        "Answer the user's question primarily using the DOCUMENT CONTEXT below. "
        "You may use minimal general knowledge only to clarify a term the document "
        "assumes the reader knows, but never to add facts the document does not support."
    )

    prompt = f"""
You are a precise, multilingual document question-answering assistant. Accuracy
matters more than fluency — never invent facts, numbers, or names.

{grounding_rule}

Target answer language: {language}

Rules:
1. Answer in {language}.
2. Base every claim strictly on the DOCUMENT CONTEXT. If two parts of the context
   conflict, point out the discrepancy instead of picking one silently.
3. If the answer is not present in the context, clearly say the information could
   not be found in the uploaded document — do not guess.
4. When you use information that came from an image, chart, diagram, or scanned
   page (marked "Image description" or "Scanned page OCR" below), say so explicitly,
   e.g. "According to the chart on page 3...".
5. Keep the answer clear, well-structured, and easy to understand.
6. Preserve important technical terms, numbers, and proper nouns exactly as given.
7. Mention the relevant page/slide/location number(s) at the end.
8. If the user asks for an explanation, explain simply without losing accuracy.

DOCUMENT CONTEXT:

{context}

USER QUESTION:

{question}
"""

    response = client.models.generate_content(
        model=GENERATION_MODEL,
        contents=prompt,
        config=types.GenerateContentConfig(temperature=0.15),
    )

    return response.text


# =========================================================
# CHAT HISTORY
# =========================================================

if st.session_state.messages:

    st.markdown('<div class="section-heading">💬 Conversation</div>', unsafe_allow_html=True)

    for message in st.session_state.messages:
        if message["role"] == "user":
            st.markdown(f'<div style="margin:12px 0 12px 10%;padding:16px 20px;border-radius:20px 20px 5px 20px;background:linear-gradient(135deg,#7c3aed,#2563eb);color:white;">👤 {message["content"]}</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div style="margin:12px 10% 12px 0;padding:18px 20px;border-radius:20px 20px 20px 5px;background:rgba(30,41,59,.88);color:#e2e8f0;border:1px solid rgba(255,255,255,.08);">🤖 {message["content"]}</div>', unsafe_allow_html=True)


# =========================================================
# QUESTION INPUT
# =========================================================

if st.session_state.processed:

    question = st.chat_input("Ask a question about your document...")

    if question:

        st.session_state.messages.append({"role": "user", "content": question})

        with st.chat_message("user"):
            st.markdown(question)

        with st.chat_message("assistant"):

            with st.spinner("Searching the document..."):
                results = search_document(
                    question, st.session_state.chunks, st.session_state.index, top_k
                )

            with st.spinner("Generating answer with Gemini..."):
                answer = generate_answer(
                    question, results, LANGUAGES[selected_language], strict=strict_grounding
                )

            st.markdown(answer)

            st.divider()
            st.caption("📑 Sources")

            for result in results:
                badge_class = {
                    "text": "type-text",
                    "image": "type-image",
                    "ocr": "type-ocr",
                }.get(result["type"], "type-text")
                badge_label = {
                    "text": "TEXT",
                    "image": "IMAGE",
                    "ocr": "SCANNED/OCR",
                }.get(result["type"], "TEXT")

                st.markdown(
                    f'<div class="source-card">'
                    f'<span class="type-badge {badge_class}">{badge_label}</span>'
                    f'<b>{result["page"]}</b> · relevance {result["score"]:.2f}'
                    f'</div>',
                    unsafe_allow_html=True,
                )

        st.session_state.messages.append({"role": "assistant", "content": answer})


# =========================================================
# EXTRA FEATURES
# =========================================================

if st.session_state.processed:

    st.divider()
    st.subheader("✨ Extra Tools")

    col1, col2 = st.columns(2)

    with col1:
        if st.button("📝 Generate Summary"):
            results = search_document(
                "Give an overall summary of this document, including key points from any "
                "images, charts, or diagrams.",
                st.session_state.chunks,
                st.session_state.index,
                top_k=min(10, len(st.session_state.chunks)),
            )
            summary = generate_answer(
                "Give an overall, well-structured summary of the available document content, "
                "explicitly mentioning insights drawn from any images or diagrams.",
                results,
                LANGUAGES[selected_language],
                strict=strict_grounding,
            )
            st.write(summary)

    with col2:
        if st.button("❓ Generate MCQs"):
            results = search_document(
                "Important concepts, facts, and figures in this document, including anything "
                "shown in images, charts, or diagrams.",
                st.session_state.chunks,
                st.session_state.index,
                top_k=min(10, len(st.session_state.chunks)),
            )
            mcqs = generate_answer(
                "Create 5 multiple-choice questions from the retrieved document content "
                "(including any image-derived facts). Give four options each, mark the "
                "correct answer clearly, and cite the page/location for each question.",
                results,
                LANGUAGES[selected_language],
                strict=strict_grounding,
            )
            st.write(mcqs)


# =========================================================
# FOOTER
# =========================================================

st.markdown("""
<div class="footer">
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━<br>
📚 <b>LinguaDoc AI</b><br>
Multilingual • Multimodal • Document-Aware<br>
English • தமிழ் • ಕನ್ನಡ • हिन्दी • తెలుగు • മലയാളం<br>
Built with Python • Streamlit • FAISS • Sentence Transformers • PyMuPDF • Gemini Vision
</div>
""", unsafe_allow_html=True)
